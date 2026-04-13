from __future__ import annotations

import base64
import io
import json
import os
import re
import uuid
from datetime import datetime, timezone
from http import HTTPStatus
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from typing import Any

try:
    import psycopg
except ImportError:  # pragma: no cover - local dev fallback
    psycopg = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml import parse_xml
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover - dependency is installed by requirements.txt
    Presentation = None
    RGBColor = None
    MSO_LINE_DASH_STYLE = None
    MSO_CONNECTOR = None
    MSO_SHAPE = None
    MSO_ANCHOR = None
    MSO_AUTO_SIZE = None
    PP_ALIGN = None
    parse_xml = None
    Inches = None
    Pt = None


APP_ROOT = Path(__file__).resolve().parent
LOCAL_FEEDBACK_DIR = Path(os.getenv("LOCAL_FEEDBACK_DIR", "/tmp/tax-flow-chart-feedback"))
DATABASE_URL = os.getenv("DATABASE_URL", "").strip()
ALLOWED_CORS_ORIGINS = {"https://jyen2k.github.io"}

CREATE_FEEDBACK_TABLE_SQL = """
CREATE TABLE IF NOT EXISTS feedback_submissions (
    id BIGSERIAL PRIMARY KEY,
    created_at TIMESTAMPTZ NOT NULL DEFAULT NOW(),
    submitter_name TEXT NOT NULL,
    comment TEXT NOT NULL,
    diagram_png BYTEA NOT NULL,
    diagram_mime_type TEXT NOT NULL DEFAULT 'image/png',
    diagram_width INTEGER,
    diagram_height INTEGER,
    authenticated_user TEXT
)
"""


def json_response(handler: "TaxChartHandler", status: int, payload: dict[str, Any]) -> None:
    body = json.dumps(payload).encode("utf-8")
    handler.send_response(status)
    add_cors_headers(handler)
    handler.send_header("Content-Type", "application/json; charset=utf-8")
    handler.send_header("Content-Length", str(len(body)))
    handler.end_headers()
    handler.wfile.write(body)


def add_cors_headers(handler: "TaxChartHandler") -> None:
    origin = handler.headers.get("Origin", "")
    if origin not in ALLOWED_CORS_ORIGINS:
        return
    handler.send_header("Access-Control-Allow-Origin", origin)
    handler.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
    handler.send_header("Access-Control-Allow-Headers", "Content-Type")
    handler.send_header("Vary", "Origin")


def ensure_feedback_table() -> None:
    if not DATABASE_URL:
        return
    if psycopg is None:
        raise RuntimeError("DATABASE_URL is set, but psycopg is not installed.")

    with psycopg.connect(DATABASE_URL) as connection:
        with connection.cursor() as cursor:
            cursor.execute(CREATE_FEEDBACK_TABLE_SQL)
        connection.commit()


def store_feedback_in_database(
    *,
    submitter_name: str,
    comment: str,
    image_bytes: bytes,
    image_mime_type: str,
    image_width: int | None,
    image_height: int | None,
    authenticated_user: str | None,
) -> dict[str, Any]:
    ensure_feedback_table()
    assert psycopg is not None

    with psycopg.connect(DATABASE_URL) as connection:
        with connection.cursor() as cursor:
            cursor.execute(
                """
                INSERT INTO feedback_submissions (
                    submitter_name,
                    comment,
                    diagram_png,
                    diagram_mime_type,
                    diagram_width,
                    diagram_height,
                    authenticated_user
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s)
                RETURNING id, created_at
                """,
                (
                    submitter_name,
                    comment,
                    image_bytes,
                    image_mime_type,
                    image_width,
                    image_height,
                    authenticated_user,
                ),
            )
            record_id, created_at = cursor.fetchone()
        connection.commit()

    return {
        "storage": "database",
        "id": record_id,
        "created_at": created_at.isoformat() if created_at else None,
    }


def store_feedback_locally(
    *,
    submitter_name: str,
    comment: str,
    image_bytes: bytes,
    image_mime_type: str,
    image_width: int | None,
    image_height: int | None,
    authenticated_user: str | None,
) -> dict[str, Any]:
    LOCAL_FEEDBACK_DIR.mkdir(parents=True, exist_ok=True)
    submission_id = uuid.uuid4().hex
    timestamp = datetime.now(timezone.utc).isoformat()
    image_path = LOCAL_FEEDBACK_DIR / f"{submission_id}.png"
    metadata_path = LOCAL_FEEDBACK_DIR / f"{submission_id}.json"

    image_path.write_bytes(image_bytes)
    metadata_path.write_text(
        json.dumps(
            {
                "id": submission_id,
                "created_at": timestamp,
                "submitter_name": submitter_name,
                "comment": comment,
                "diagram_mime_type": image_mime_type,
                "diagram_width": image_width,
                "diagram_height": image_height,
                "authenticated_user": authenticated_user,
                "image_path": str(image_path),
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    return {
        "storage": "local",
        "id": submission_id,
        "created_at": timestamp,
        "image_path": str(image_path),
    }


def list_feedback_from_database(limit: int = 25) -> list[dict[str, Any]]:
    ensure_feedback_table()
    assert psycopg is not None

    with psycopg.connect(DATABASE_URL) as connection:
        with connection.cursor() as cursor:
            cursor.execute(
                """
                SELECT
                    id,
                    created_at,
                    submitter_name,
                    comment,
                    diagram_mime_type,
                    encode(diagram_png, 'base64') AS diagram_base64
                FROM feedback_submissions
                ORDER BY created_at DESC
                LIMIT %s
                """,
                (limit,),
            )
            rows = cursor.fetchall()

    entries = []
    for row in rows:
        record_id, created_at, submitter_name, comment, diagram_mime_type, diagram_base64 = row
        entries.append(
            {
                "id": record_id,
                "created_at": created_at.isoformat() if created_at else None,
                "submitter_name": submitter_name,
                "comment": comment,
                "snapshot_data_url": f"data:{diagram_mime_type};base64,{diagram_base64}",
            }
        )
    return entries


def list_feedback_locally(limit: int = 25) -> list[dict[str, Any]]:
    if not LOCAL_FEEDBACK_DIR.exists():
        return []

    entries = []
    for metadata_path in LOCAL_FEEDBACK_DIR.glob("*.json"):
        try:
            payload = json.loads(metadata_path.read_text(encoding="utf-8"))
            image_path = Path(payload.get("image_path", ""))
            if not image_path.exists():
                continue
            image_base64 = base64.b64encode(image_path.read_bytes()).decode("ascii")
            entries.append(
                {
                    "id": payload.get("id"),
                    "created_at": payload.get("created_at"),
                    "submitter_name": payload.get("submitter_name"),
                    "comment": payload.get("comment"),
                    "snapshot_data_url": f"data:{payload.get('diagram_mime_type', 'image/png')};base64,{image_base64}",
                }
            )
        except Exception as error:  # pragma: no cover - local corruption safeguard
            print(f"Skipping unreadable feedback file {metadata_path}: {error}")

    entries.sort(key=lambda entry: entry.get("created_at") or "", reverse=True)
    return entries[:limit]


def delete_feedback_from_database(entry_id: str) -> bool:
    ensure_feedback_table()
    assert psycopg is not None

    with psycopg.connect(DATABASE_URL) as connection:
        with connection.cursor() as cursor:
            cursor.execute("DELETE FROM feedback_submissions WHERE id = %s RETURNING id", (entry_id,))
            deleted = cursor.fetchone()
        connection.commit()

    return bool(deleted)


def delete_feedback_locally(entry_id: str) -> bool:
    if not re.fullmatch(r"[a-f0-9]{32}", entry_id):
        return False

    metadata_path = LOCAL_FEEDBACK_DIR / f"{entry_id}.json"
    image_path = LOCAL_FEEDBACK_DIR / f"{entry_id}.png"
    deleted = False

    if metadata_path.exists():
        metadata_path.unlink()
        deleted = True
    if image_path.exists():
        image_path.unlink()
        deleted = True

    return deleted


BOX_WIDTH = 120
BOX_HEIGHT = 80


def pptx_response(handler: "TaxChartHandler", content: bytes) -> None:
    handler.send_response(HTTPStatus.OK)
    add_cors_headers(handler)
    handler.send_header(
        "Content-Type",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    handler.send_header("Content-Disposition", 'attachment; filename="tax-structure-diagram.pptx"')
    handler.send_header("Content-Length", str(len(content)))
    handler.end_headers()
    handler.wfile.write(content)


def srgb(hex_value: str) -> str:
    return hex_value.replace("#", "").upper()


def node_lookup(payload: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {node.get("id"): node for node in payload.get("nodes", []) if node.get("id")}


def ownership_render_nodes(edge: dict[str, Any], nodes: dict[str, dict[str, Any]]) -> tuple[dict[str, Any], dict[str, Any]] | None:
    from_node = nodes.get(edge.get("from"))
    to_node = nodes.get(edge.get("to"))
    if not from_node or not to_node:
        return None
    if edge.get("preserveDirection"):
        return from_node, to_node

    from_center_y = from_node.get("y", 0) + BOX_HEIGHT / 2
    to_center_y = to_node.get("y", 0) + BOX_HEIGHT / 2
    if from_center_y < to_center_y:
        return from_node, to_node
    if to_center_y < from_center_y:
        return to_node, from_node
    return (from_node, to_node) if from_node.get("x", 0) <= to_node.get("x", 0) else (to_node, from_node)


def diagram_bounds(payload: dict[str, Any]) -> tuple[float, float, float, float]:
    nodes = payload.get("nodes", [])
    if not nodes:
        return 0, 0, 1400, 900
    min_x = min(float(node.get("x", 0)) for node in nodes)
    min_y = min(float(node.get("y", 0)) for node in nodes)
    max_x = max(float(node.get("x", 0)) + BOX_WIDTH for node in nodes)
    max_y = max(
        float(node.get("y", 0))
        + (BOX_HEIGHT + 38 if node.get("type") == "individual" else BOX_HEIGHT)
        for node in nodes
    )
    return min_x - 80, min_y - 80, max_x + 80, max_y + 120


PPTX_BASE_INCHES_PER_UNIT = 1 / 80
PPTX_SLIDE_MARGIN_INCHES = 0.45
PPTX_MIN_SLIDE_WIDTH_INCHES = 7.5
PPTX_MIN_SLIDE_HEIGHT_INCHES = 5.625
PPTX_MAX_SLIDE_DIMENSION_INCHES = 24


def pptx_rgb(hex_value: str) -> Any:
    value = srgb(hex_value)
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


class EditablePptxDiagramBuilder:
    """Builds a PowerPoint file from real editable PowerPoint shapes."""

    def __init__(self, payload: dict[str, Any]) -> None:
        if Presentation is None:
            raise RuntimeError("PowerPoint export dependency is not installed.")

        self.payload = payload
        self.nodes = node_lookup(payload)
        min_x, min_y, max_x, max_y = self.diagram_bounds()
        self.min_x = min_x
        self.min_y = min_y
        self.diagram_width = max(1, max_x - min_x)
        self.diagram_height = max(1, max_y - min_y)

        raw_width = self.diagram_width * PPTX_BASE_INCHES_PER_UNIT + PPTX_SLIDE_MARGIN_INCHES * 2
        raw_height = self.diagram_height * PPTX_BASE_INCHES_PER_UNIT + PPTX_SLIDE_MARGIN_INCHES * 2
        self.slide_width_inches = max(
            PPTX_MIN_SLIDE_WIDTH_INCHES,
            min(PPTX_MAX_SLIDE_DIMENSION_INCHES, raw_width),
        )
        self.slide_height_inches = max(
            PPTX_MIN_SLIDE_HEIGHT_INCHES,
            min(PPTX_MAX_SLIDE_DIMENSION_INCHES, raw_height),
        )
        self.scale = min(
            (self.slide_width_inches - 2 * PPTX_SLIDE_MARGIN_INCHES) / self.diagram_width,
            (self.slide_height_inches - 2 * PPTX_SLIDE_MARGIN_INCHES) / self.diagram_height,
        )

        self.presentation = Presentation()
        self.presentation.slide_width = Inches(self.slide_width_inches)
        self.presentation.slide_height = Inches(self.slide_height_inches)
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.slide.background.fill.solid()
        self.slide.background.fill.fore_color.rgb = pptx_rgb("FFFFFF")

    def diagram_bounds(self) -> tuple[float, float, float, float]:
        min_x, min_y, max_x, max_y = diagram_bounds(self.payload)
        legend = self.payload.get("transactionLegend") or {}
        if legend.get("enabled"):
            max_x = max(max_x, float(legend.get("x", 70)) + 340)
            max_y = max(max_y, float(legend.get("y", 760)) + 96)
        return min_x, min_y, max_x, max_y

    def x(self, value: float) -> Any:
        return Inches(PPTX_SLIDE_MARGIN_INCHES + (value - self.min_x) * self.scale)

    def y(self, value: float) -> Any:
        return Inches(PPTX_SLIDE_MARGIN_INCHES + (value - self.min_y) * self.scale)

    def w(self, value: float) -> Any:
        return Inches(max(0.01, value * self.scale))

    def h(self, value: float) -> Any:
        return Inches(max(0.01, value * self.scale))

    def add_shape(
        self,
        shape_type: Any,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        fill: str | None = "FFFDF9",
        line: str | None = "514236",
        dashed: bool = False,
        line_width_pt: float = 1.25,
    ) -> Any:
        shape = self.slide.shapes.add_shape(shape_type, self.x(x), self.y(y), self.w(width), self.h(height))
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = pptx_rgb(fill)
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = pptx_rgb(line)
            shape.line.width = Pt(line_width_pt)
            if dashed:
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            shape.line.fill.background()
        return shape

    def add_line(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        color: str = "1E1A17",
        width_pt: float = 1.25,
        dashed: bool = False,
        arrow_start: bool = False,
        arrow_end: bool = False,
        connector_type: Any | None = None,
    ) -> Any:
        connector = self.slide.shapes.add_connector(
            connector_type or MSO_CONNECTOR.STRAIGHT,
            self.x(x1),
            self.y(y1),
            self.x(x2),
            self.y(y2),
        )
        connector.line.color.rgb = pptx_rgb(color)
        connector.line.width = Pt(width_pt)
        if dashed:
            connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if arrow_start or arrow_end:
            self.add_arrowheads(connector, arrow_start=arrow_start, arrow_end=arrow_end)
        return connector

    def add_arrowheads(self, connector: Any, *, arrow_start: bool, arrow_end: bool) -> None:
        if parse_xml is None:
            return
        line_properties = connector.line._get_or_add_ln()
        if arrow_start:
            line_properties.append(
                parse_xml(
                    '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>',
                ),
            )
        if arrow_end:
            line_properties.append(
                parse_xml(
                    '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle"/>',
                ),
            )

    def add_textbox(
        self,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        *,
        font_size: float = 9,
        color: str = "23180F",
        bold: bool = False,
        align: Any | None = None,
        anchor: Any | None = None,
    ) -> Any:
        textbox = self.slide.shapes.add_textbox(self.x(x), self.y(y), self.w(width), self.h(height))
        textbox.fill.background()
        textbox.line.fill.background()
        frame = textbox.text_frame
        frame.clear()
        frame.margin_left = Inches(0.03)
        frame.margin_right = Inches(0.03)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.vertical_anchor = anchor or MSO_ANCHOR.MIDDLE

        lines = str(text or "").splitlines() or [""]
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = align or PP_ALIGN.CENTER
            run = paragraph.add_run()
            run.text = line
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = pptx_rgb(color)
        return textbox

    def wrap_node_label(self, node: dict[str, Any]) -> list[str]:
        label = str(node.get("label") or "").strip()
        words = label.split()
        if not words:
            lines = [""]
        else:
            max_chars = {"partnership": 11, "trust": 12, "individual": 14}.get(node.get("type"), 15)
            lines = []
            current = ""
            for word in words:
                candidate = f"{current} {word}".strip()
                if not current or len(candidate) <= max_chars:
                    current = candidate
                else:
                    lines.append(current)
                    current = word
            if current:
                lines.append(current)

        jurisdiction = self.jurisdiction_text(node)
        if jurisdiction:
            lines.append(jurisdiction)
        return lines

    def jurisdiction_text(self, node: dict[str, Any]) -> str:
        custom = str(node.get("jurisdictionCustom") or "").strip()
        jurisdiction = str(node.get("jurisdiction") or "").strip()
        if node.get("jurisdictionMode") == "custom" and custom:
            return custom
        if len(jurisdiction) == 2 and jurisdiction.isalpha():
            return "".join(chr(127397 + ord(letter.upper())) for letter in jurisdiction)
        return jurisdiction

    def label_text(self, node: dict[str, Any]) -> str:
        return "\n".join(self.wrap_node_label(node))

    def ownership_start_y(self, node: dict[str, Any]) -> float:
        if node.get("type") != "individual":
            return float(node.get("y", 0)) + BOX_HEIGHT
        label_line_count = max(1, len(self.wrap_node_label(node)))
        return float(node.get("y", 0)) + BOX_HEIGHT + 8 + label_line_count * 12

    def render_nodes(self) -> None:
        for node in self.payload.get("nodes", []):
            x = float(node.get("x", 0))
            y = float(node.get("y", 0))
            node_type = node.get("type", "corporation")
            fill = "E7E1D6" if node.get("fill") == "shaded" else "FFFDF9"
            dashed = node.get("lineStyle") == "dashed"
            label = self.label_text(node)

            if node_type == "individual":
                self.render_individual(node, x, y, label)
            elif node_type == "trust":
                self.add_shape(MSO_SHAPE.OVAL, x, y, BOX_WIDTH, BOX_HEIGHT, fill=fill, dashed=dashed)
                self.add_textbox(label, x + 5, y, BOX_WIDTH - 10, BOX_HEIGHT, font_size=9)
            elif node_type == "partnership":
                self.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, BOX_WIDTH, BOX_HEIGHT, fill=fill, dashed=dashed)
                self.add_textbox(label, x + 10, y + 25, BOX_WIDTH - 20, BOX_HEIGHT - 30, font_size=8.5)
            else:
                self.add_shape(MSO_SHAPE.RECTANGLE, x, y, BOX_WIDTH, BOX_HEIGHT, fill=fill, dashed=dashed)
                if node_type in {"dreg", "hybrid"}:
                    inner_fill = "E7E1D6" if node.get("innerFill") == "shaded" else None
                    self.add_shape(
                        MSO_SHAPE.OVAL,
                        x,
                        y,
                        BOX_WIDTH,
                        BOX_HEIGHT,
                        fill=inner_fill,
                        dashed=node.get("innerLineStyle") == "dashed",
                    )
                elif node_type == "hybrid-partnership":
                    inner_dashed = node.get("innerLineStyle") == "dashed"
                    self.add_line(x, y + BOX_HEIGHT, x + BOX_WIDTH / 2, y, dashed=inner_dashed)
                    self.add_line(x + BOX_WIDTH, y + BOX_HEIGHT, x + BOX_WIDTH / 2, y, dashed=inner_dashed)
                elif node_type == "reverse-hybrid":
                    inner_dashed = node.get("innerLineStyle") == "dashed"
                    self.add_line(x, y, x + BOX_WIDTH / 2, y + BOX_HEIGHT, dashed=inner_dashed)
                    self.add_line(x + BOX_WIDTH, y, x + BOX_WIDTH / 2, y + BOX_HEIGHT, dashed=inner_dashed)
                self.add_textbox(label, x + 5, y, BOX_WIDTH - 10, BOX_HEIGHT, font_size=9)

            if node.get("crossedOut") and node_type != "individual":
                self.add_line(x - 8, y - 8, x + BOX_WIDTH + 8, y + BOX_HEIGHT + 8, dashed=True)
                self.add_line(x + BOX_WIDTH + 8, y - 8, x - 8, y + BOX_HEIGHT + 8, dashed=True)

    def render_individual(self, node: dict[str, Any], x: float, y: float, label: str) -> None:
        centers = [BOX_WIDTH / 2]
        if node.get("multipleIndividuals"):
            centers = [BOX_WIDTH / 2 - 22, BOX_WIDTH / 2 + 22, BOX_WIDTH / 2]
        for index, center in enumerate(centers):
            top = y + (10 if node.get("multipleIndividuals") and index < 2 else 2)
            cx = x + center
            self.add_shape(MSO_SHAPE.OVAL, cx - 9, top, 18, 18, fill=None, line="514236", line_width_pt=0.9)
            self.add_line(cx, top + 18, cx, top + 46, width_pt=0.9)
            self.add_line(cx - 18, top + 30, cx + 18, top + 30, width_pt=0.9)
            self.add_line(cx, top + 46, cx - 16, top + 68, width_pt=0.9)
            self.add_line(cx, top + 46, cx + 16, top + 68, width_pt=0.9)
        self.add_textbox(label, x, y + BOX_HEIGHT + 2, BOX_WIDTH, 36, font_size=9)

    def render_edges(self, kind: str) -> None:
        for edge in self.payload.get("edges", []):
            if edge.get("kind") != kind:
                continue
            if kind == "ownership":
                self.render_ownership(edge)
            elif kind == "transaction":
                self.render_transaction(edge)

    def edge_color(self, edge: dict[str, Any]) -> str:
        return {"red": "B0392F", "blue": "245EA8"}.get(edge.get("color"), "1E1A17")

    def render_ownership(self, edge: dict[str, Any]) -> None:
        nodes = ownership_render_nodes(edge, self.nodes)
        if not nodes:
            return
        parent, child = nodes
        start_x = float(parent.get("x", 0)) + BOX_WIDTH / 2
        start_y = self.ownership_start_y(parent)
        end_x = float(child.get("x", 0)) + BOX_WIDTH / 2
        end_y = float(child.get("y", 0))
        mid_y = start_y + (end_y - start_y) / 2
        color = self.edge_color(edge)
        dashed = edge.get("lineStyle") == "dashed"

        self.add_line(start_x, start_y, start_x, mid_y, color=color, dashed=dashed)
        self.add_line(start_x, mid_y, end_x, mid_y, color=color, dashed=dashed)
        self.add_line(end_x, mid_y, end_x, end_y, color=color, dashed=dashed)

        label = "\n".join(str(value) for value in [edge.get("percent"), edge.get("label")] if value)
        if not label:
            return
        side = -1 if end_x <= start_x else 1
        label_width = 60
        label_x = end_x + side * 14 - (label_width if side < 0 else 0)
        self.add_textbox(
            label,
            label_x,
            (mid_y + end_y) / 2 - 12,
            label_width,
            26,
            font_size=7.5,
            color=color,
            align=PP_ALIGN.RIGHT if side < 0 else PP_ALIGN.LEFT,
        )

    def render_transaction(self, edge: dict[str, Any]) -> None:
        from_node = self.nodes.get(edge.get("from"))
        to_node = self.nodes.get(edge.get("to"))
        if not from_node or not to_node:
            return

        from_left = float(from_node.get("x", 0)) > float(to_node.get("x", 0))
        fx = float(from_node.get("x", 0)) + (0 if from_left else BOX_WIDTH)
        tx = float(to_node.get("x", 0)) + (BOX_WIDTH if from_left else 0)
        fy = float(from_node.get("y", 0)) + BOX_HEIGHT / 2
        ty = float(to_node.get("y", 0)) + BOX_HEIGHT / 2
        color = self.edge_color(edge)
        dashed = edge.get("lineStyle") == "dashed"
        connector_type = MSO_CONNECTOR.CURVE if edge.get("curveOffset") else MSO_CONNECTOR.STRAIGHT

        self.add_line(
            fx,
            fy,
            tx,
            ty,
            color=color,
            dashed=dashed,
            width_pt=1.6,
            arrow_end=True,
            arrow_start=bool(edge.get("bidirectional")),
            connector_type=connector_type,
        )
        self.add_textbox(
            edge.get("label") or "Transaction",
            (fx + tx) / 2 - 55,
            (fy + ty) / 2 - 13,
            110,
            26,
            font_size=7.5,
            color=color,
        )

    def render_legend(self) -> None:
        legend = self.payload.get("transactionLegend") or {}
        if not legend.get("enabled"):
            return
        arrow_text = str(legend.get("arrowEndText") or "").strip()
        non_arrow_text = str(legend.get("nonArrowEndText") or "").strip()
        if not arrow_text and not non_arrow_text:
            return
        x = float(legend.get("x", 70))
        y = float(legend.get("y", 760))
        width = 340
        height = 96
        line_y = y + 64
        start_x = x + 20
        end_x = x + width - 20

        self.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height, fill="FFFDF9", line="D9C7B7", line_width_pt=0.9)
        self.add_textbox("Legend", x + 16, y + 8, 120, 24, font_size=9, align=PP_ALIGN.LEFT)
        self.add_line(start_x, line_y, end_x, line_y, arrow_end=True, width_pt=1.4)
        if non_arrow_text:
            self.add_textbox(non_arrow_text, start_x, line_y - 26, 120, 22, font_size=7.5, align=PP_ALIGN.LEFT)
        if arrow_text:
            self.add_textbox(arrow_text, end_x - 120, line_y - 26, 120, 22, font_size=7.5, align=PP_ALIGN.RIGHT)

    def build(self) -> bytes:
        self.render_edges("ownership")
        self.render_nodes()
        self.render_edges("transaction")
        self.render_legend()
        buffer = io.BytesIO()
        self.presentation.save(buffer)
        return buffer.getvalue()


def build_pptx(payload: dict[str, Any]) -> bytes:
    return EditablePptxDiagramBuilder(payload).build()

class TaxChartHandler(SimpleHTTPRequestHandler):
    def __init__(self, *args: Any, **kwargs: Any) -> None:
        super().__init__(*args, directory=str(APP_ROOT), **kwargs)

    def do_OPTIONS(self) -> None:
        if self.path.rstrip("/") != "/api/export/pptx":
            self.send_error(HTTPStatus.NOT_FOUND, "Unknown API endpoint.")
            return

        self.send_response(HTTPStatus.NO_CONTENT)
        add_cors_headers(self)
        self.send_header("Content-Length", "0")
        self.end_headers()

    def do_GET(self) -> None:
        if self.path.rstrip("/") == "/api/feedback":
            try:
                entries = list_feedback_from_database() if DATABASE_URL else list_feedback_locally()
                json_response(self, HTTPStatus.OK, {"ok": True, "entries": entries})
            except Exception as error:  # pragma: no cover - runtime safeguard
                print(f"Feedback inbox load failed: {error}")
                json_response(
                    self,
                    HTTPStatus.INTERNAL_SERVER_ERROR,
                    {"ok": False, "error": "Could not load feedback right now."},
                )
            return

        if self.path in {"", "/"}:
            self.path = "/index.html"
        return super().do_GET()

    def do_POST(self) -> None:
        if self.path.rstrip("/") == "/api/export/pptx":
            try:
                self.handle_pptx_export()
            except ValueError as error:
                json_response(self, HTTPStatus.BAD_REQUEST, {"ok": False, "error": str(error)})
            except Exception as error:  # pragma: no cover - runtime safeguard
                print(f"PPTX export failed: {error}")
                json_response(
                    self,
                    HTTPStatus.INTERNAL_SERVER_ERROR,
                    {"ok": False, "error": "Could not export PowerPoint right now."},
                )
            return

        if self.path.rstrip("/") != "/api/feedback":
            self.send_error(HTTPStatus.NOT_FOUND, "Unknown API endpoint.")
            return

        try:
            self.handle_feedback_submission()
        except ValueError as error:
            json_response(self, HTTPStatus.BAD_REQUEST, {"ok": False, "error": str(error)})
        except Exception as error:  # pragma: no cover - runtime safeguard
            print(f"Feedback submission failed: {error}")
            json_response(
                self,
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"ok": False, "error": "Could not store feedback right now."},
            )

    def handle_pptx_export(self) -> None:
        content_length = int(self.headers.get("Content-Length", "0") or "0")
        if content_length <= 0:
            raise ValueError("Missing diagram payload.")
        if content_length > 5_000_000:
            raise ValueError("Diagram payload is too large.")

        raw_body = self.rfile.read(content_length)
        payload = json.loads(raw_body.decode("utf-8"))
        if not isinstance(payload.get("nodes"), list):
            raise ValueError("Diagram payload is missing nodes.")

        pptx_response(self, build_pptx(payload))

    def do_DELETE(self) -> None:
        if self.path.rstrip("/").split("?")[0] != "/api/feedback":
            self.send_error(HTTPStatus.NOT_FOUND, "Unknown API endpoint.")
            return

        try:
            self.handle_feedback_delete()
        except ValueError as error:
            json_response(self, HTTPStatus.BAD_REQUEST, {"ok": False, "error": str(error)})
        except Exception as error:  # pragma: no cover - runtime safeguard
            print(f"Feedback delete failed: {error}")
            json_response(
                self,
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"ok": False, "error": "Could not delete feedback right now."},
            )

    def handle_feedback_delete(self) -> None:
        request_path = self.path.split("?", 1)
        query = request_path[1] if len(request_path) > 1 else ""
        params = {}
        for pair in query.split("&"):
            if "=" not in pair:
                continue
            key, value = pair.split("=", 1)
            params[key] = value

        entry_id = params.get("id", "").strip()
        if not entry_id:
            raise ValueError("Feedback id is required.")

        deleted = (
            delete_feedback_from_database(entry_id)
            if DATABASE_URL
            else delete_feedback_locally(entry_id)
        )

        if not deleted:
            json_response(self, HTTPStatus.NOT_FOUND, {"ok": False, "error": "Feedback item not found."})
            return

        json_response(self, HTTPStatus.OK, {"ok": True, "message": "Feedback deleted."})

    def handle_feedback_submission(self) -> None:
        content_length = int(self.headers.get("Content-Length", "0") or "0")
        if content_length <= 0:
            raise ValueError("Missing feedback payload.")
        if content_length > 20_000_000:
            raise ValueError("Feedback payload is too large.")

        raw_body = self.rfile.read(content_length)
        payload = json.loads(raw_body.decode("utf-8"))

        submitter_name = str(payload.get("name", "")).strip()
        comment = str(payload.get("comment", "")).strip()
        snapshot_base64 = str(payload.get("snapshotBase64", "")).strip()
        snapshot_mime_type = str(payload.get("snapshotMimeType", "image/png")).strip() or "image/png"
        snapshot_width = payload.get("snapshotWidth")
        snapshot_height = payload.get("snapshotHeight")
        authenticated_user = self.headers.get("X-Authenticated-User", "").strip() or None

        if not submitter_name:
            raise ValueError("Name is required.")
        if not comment:
            raise ValueError("Comment is required.")
        if not snapshot_base64:
            raise ValueError("Diagram snapshot is required.")

        try:
            image_bytes = base64.b64decode(snapshot_base64, validate=True)
        except Exception as error:
            raise ValueError("Diagram snapshot could not be decoded.") from error

        if DATABASE_URL:
            result = store_feedback_in_database(
                submitter_name=submitter_name,
                comment=comment,
                image_bytes=image_bytes,
                image_mime_type=snapshot_mime_type,
                image_width=snapshot_width,
                image_height=snapshot_height,
                authenticated_user=authenticated_user,
            )
        else:
            result = store_feedback_locally(
                submitter_name=submitter_name,
                comment=comment,
                image_bytes=image_bytes,
                image_mime_type=snapshot_mime_type,
                image_width=snapshot_width,
                image_height=snapshot_height,
                authenticated_user=authenticated_user,
            )

        json_response(
            self,
            HTTPStatus.OK,
            {
                "ok": True,
                "message": "Feedback saved.",
                **result,
            },
        )


def main() -> None:
    port = int(os.getenv("PORT", "8000"))
    server = ThreadingHTTPServer(("0.0.0.0", port), TaxChartHandler)
    print(f"Serving tax chart tool on http://0.0.0.0:{port}")
    server.serve_forever()


if __name__ == "__main__":
    main()
